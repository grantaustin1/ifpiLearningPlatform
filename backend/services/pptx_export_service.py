"""PPTX export for a course (Iter 27c).

Uses python-pptx to render each slide (title + HTML-stripped body) into a
downloadable `.pptx`. Runs sync — for large courses (200+ slides) we may
want to move to a background job later.
"""
from __future__ import annotations

import io
import re
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches, Pt


_TAG_RE = re.compile(r"<[^>]+>")


def _plain(text: str) -> str:
    if not text:
        return ""
    txt = re.sub(r"<br\s*/?>|</p>|</li>", "\n", text, flags=re.I)
    txt = _TAG_RE.sub("", txt)
    txt = (txt.replace("&nbsp;", " ").replace("&amp;", "&")
              .replace("&lt;", "<").replace("&gt;", ">")
              .replace("&quot;", '"'))
    return re.sub(r"\n{3,}", "\n\n", txt).strip()


def build_pptx(course_title: str, description: str, slides: Iterable[dict]) -> bytes:
    """`slides` is an iterable of `{title, content, slide_type, media_url}`.
    Returns the PPTX file as bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    # ── Cover slide ──────────────────────────────────────────────
    cover_layout = prs.slide_layouts[0]   # Title Slide
    cover = prs.slides.add_slide(cover_layout)
    cover.shapes.title.text = course_title or "Course"
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = (description or "IFPI Learning Platform")

    # ── Content slides ──────────────────────────────────────────
    content_layout = prs.slide_layouts[1]   # Title + Content
    for s in slides:
        title = (s.get("title") or "Untitled").strip()[:120]
        body = _plain(s.get("content") or "")[:2500]
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title

        # Body placeholder — index 1 in default template
        if len(slide.placeholders) > 1:
            body_ph = slide.placeholders[1]
            tf = body_ph.text_frame
            tf.clear()
            # Split into paragraphs
            paragraphs = [p for p in body.split("\n") if p.strip()]
            if not paragraphs:
                paragraphs = ["(no body content)"]
            for i, p in enumerate(paragraphs):
                target = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                target.text = p.strip()
                target.font.size = Pt(14)

        # Media annotation (we can't embed remote URLs safely, so we cite them)
        if s.get("media_url"):
            annotation = (
                f"\n\n[Media: {s.get('slide_type', 'MEDIA').lower()} → "
                f"{s['media_url']}]"
            )
            if len(slide.placeholders) > 1:
                p = slide.placeholders[1].text_frame.add_paragraph()
                p.text = annotation.strip()
                p.font.size = Pt(10)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
